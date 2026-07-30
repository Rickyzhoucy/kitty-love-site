"""联网与文档生成两组新工具。"""

import zipfile
from io import BytesIO

import pytest
from sqlalchemy import select

from app.agent_tasks import describe_step
from app.agent_tools import build_domain_tools
from app.agents.roles import AgentRole, filter_tools
from app.config import Settings
from app.doc_tools import build_document_tools, persist_document
from app.document_builder import DocumentSpecError, build_document
from app.models import Attachment, User
from app.web_search import (
    BochaSearchProvider,
    BraveSearchProvider,
    TavilySearchProvider,
    build_search_provider,
)
from app.web_tools import build_web_tools, guard_url

# ---- 文档生成 ----


def _open(content: bytes) -> zipfile.ZipFile:
    """三种 Office 格式都是 ZIP。能打开就说明产出的是合法容器。"""
    return zipfile.ZipFile(BytesIO(content))


def test_docx_covers_every_block_type():
    built = build_document(
        "docx",
        "周报",
        [
            {"type": "heading", "level": 1, "text": "本周进展"},
            {"type": "paragraph", "text": "一切顺利。"},
            {"type": "bullets", "items": ["修了缩放", "打磨了菜单"]},
            {"type": "numbers", "items": ["第一", "第二"]},
            {"type": "quote", "text": "慢就是快。"},
            {"type": "table", "rows": [["项目", "状态"], ["对话本", "已上线"]]},
        ],
    )
    assert built.filename == "周报.docx"
    assert "word/document.xml" in _open(built.content).namelist()


def test_pptx_builds_slides_with_notes():
    built = build_document(
        "pptx",
        "季度汇报",
        [
            {"title": "结论", "bullets": ["做完了", "还能更好"], "notes": "别念稿"},
            {"title": "只有标题的一页"},
        ],
    )
    names = _open(built.content).namelist()
    # 封面 + 两页
    assert sum(1 for name in names if name.startswith("ppt/slides/slide")) == 3
    assert any(name.startswith("ppt/notesSlides/") for name in names)


def test_subtitle_only_block_does_not_leave_a_blank_slide():
    """模型很自然会把副标题写成第一个块。它已经被封面用掉了，
    不跳过就会多出一页彻底空白的幻灯片。"""
    from pptx import Presentation

    built = build_document(
        "pptx",
        "年度回顾",
        [{"subtitle": "一份小结"}, {"title": "正文", "bullets": ["一"]}],
    )
    slides = Presentation(BytesIO(built.content)).slides
    assert len(slides) == 2
    assert all(
        slide.shapes.title is not None and slide.shapes.title.text
        for slide in slides
    )


def test_xlsx_sanitises_sheet_names():
    """Excel 不允许工作表名里出现 []:*?/\\，带进去会存不出来。"""
    built = build_document(
        "xlsx",
        "数据",
        [{"sheet": "一月/二月[初]", "rows": [["名称", "数量"], ["猫粮", "3"]]}],
    )
    assert "xl/workbook.xml" in _open(built.content).namelist()


@pytest.mark.parametrize(
    ("kind", "blocks", "hint"),
    [
        ("pdf", [{"text": "x"}], "不支持"),
        ("docx", [], "非空"),
        ("docx", "不是数组", "非空"),
        ("docx", ["不是对象"], "对象"),
        ("xlsx", [{"rows": "不是二维数组"}], "二维数组"),
    ],
)
def test_bad_spec_raises_a_message_the_model_can_act_on(kind, blocks, hint):
    with pytest.raises(DocumentSpecError, match=hint):
        build_document(kind, "标题", blocks)


def test_oversized_spec_is_rejected():
    with pytest.raises(DocumentSpecError, match="最多"):
        build_document("docx", "长文", [{"text": "x"}] * 500)


async def test_generated_document_becomes_a_downloadable_attachment(session_maker):
    class FakeStorage:
        def __init__(self):
            self.written: list[tuple[str, str, int]] = []

        def build_object_key(self, owner_id, filename):
            return f"{owner_id}/generated/{filename}"

        async def put_bytes(self, bucket, object_key, content, content_type):
            self.written.append((bucket, object_key, len(content)))

    storage = FakeStorage()
    settings = Settings(
        database_url="sqlite+aiosqlite://",
        session_secret="x" * 32,
    )
    built = build_document("docx", "备忘", [{"text": "记一下"}])
    async with session_maker() as db:
        user_id = await db.scalar(select(User.id))
        attachment = await persist_document(
            db, storage, settings, user_id, built.filename,
            built.content_type, built.content,
        )

    assert storage.written and storage.written[0][2] == len(built.content)
    async with session_maker() as db:
        stored = await db.get(Attachment, attachment.id)
    assert stored.owner_id == user_id
    assert stored.filename == "备忘.docx"
    # 自己生成的文件不需要再解析一遍取文本
    assert stored.parse_status == "ready"


# ---- 联网 ----


def _resolve_to(monkeypatch, address: str) -> None:
    """把 DNS 钉死。不这么做，这些用例就会依赖真实网络而变得不稳定，
    而且在断网的 CI 上「拒绝访问」会因为解析失败而假阳性通过。"""
    import socket as socket_module

    monkeypatch.setattr(
        socket_module,
        "getaddrinfo",
        lambda *args, **kwargs: [(0, 0, 0, "", (address, 0))],
    )


@pytest.mark.parametrize(
    "url",
    [
        "file:///etc/passwd",
        "ftp://example.com/x",
        "not-a-url",
        "http://",
    ],
)
async def test_non_http_targets_are_refused_before_any_lookup(url):
    with pytest.raises(ValueError):
        await guard_url(url)


@pytest.mark.parametrize(
    "address",
    [
        "127.0.0.1",       # 回环
        "10.0.0.5",        # 私网
        "192.168.1.1",     # 私网
        "169.254.169.254", # 云厂商元数据服务——这条是最要命的
        "::1",             # IPv6 回环
        "fd00::1",         # IPv6 私网
    ],
)
async def test_hosts_resolving_inward_are_refused(monkeypatch, address):
    """模型给的 URL 可能来自它刚读到的网页，不设防就是一个 SSRF。

    判定必须在**解析之后**：`localtest.me` 这种域名看着是公网的，
    解析出来却是 127.0.0.1，只看字面量会被绕过。
    """
    _resolve_to(monkeypatch, address)
    with pytest.raises(ValueError, match="内网或本机"):
        await guard_url("https://looks-legit.example.com/x")


async def test_public_hosts_are_allowed(monkeypatch):
    _resolve_to(monkeypatch, "93.184.216.34")
    assert await guard_url("https://example.com/a?b=1") == "https://example.com/a?b=1"


def test_search_tool_is_not_registered_without_a_key():
    """没配 key 就不该注册这个工具——注册了再在调用时报错更糟。"""
    names = {getattr(item, "name", "") for item in build_web_tools(None)}
    assert "web_search" not in names
    assert "web_read" in names


def test_provider_selection_is_config_driven():
    base = {"database_url": "sqlite+aiosqlite://", "session_secret": "x" * 32}
    assert build_search_provider(Settings(**base, web_search_api_key="")) is None
    assert isinstance(
        build_search_provider(
            Settings(**base, web_search_api_key="k", web_search_provider="bocha")
        ),
        BochaSearchProvider,
    )
    assert isinstance(
        build_search_provider(
            Settings(**base, web_search_api_key="k", web_search_provider="tavily")
        ),
        TavilySearchProvider,
    )
    assert isinstance(
        build_search_provider(
            Settings(**base, web_search_api_key="k", web_search_provider="Brave")
        ),
        BraveSearchProvider,
    )
    # 拼错提供方名字不该悄悄退回默认的那家——那会把查询词发给意料之外的服务。
    assert build_search_provider(
        Settings(**base, web_search_api_key="k", web_search_provider="typo")
    ) is None


# ---- 权限与分级 ----


def test_pet_cannot_autonomously_search_or_generate_files(session_maker):
    """联网要花钱且会把内容发到站外，生成文档是用户要的产物——都不该自主执行。"""
    tools = [
        *build_domain_tools(session_maker),
        *build_web_tools(BochaSearchProvider("k", 5.0)),
        *build_document_tools(session_maker),
    ]
    cognition = {getattr(t, "name", "") for t in filter_tools(AgentRole.COGNITION, tools)}
    assert not cognition & {"web_search", "web_read", "create_document"}
    conversation = {
        getattr(t, "name", "") for t in filter_tools(AgentRole.CONVERSATION, tools)
    }
    assert {"web_search", "web_read", "create_document"} <= conversation


@pytest.mark.parametrize(
    ("tool_name", "capability", "risk", "external"),
    [
        ("web_search", "web.search", "low", True),
        ("web_read", "web.read", "low", True),
        ("create_document", "site.document", "low", False),
    ],
)
def test_new_tools_are_classified(tool_name, capability, risk, external):
    step = describe_step(tool_name, {})
    assert step.capability == capability
    assert step.risk_level == risk
    assert step.external is external


# ---- 生成的文件要挂回消息 ----


def test_only_create_document_contributes_an_attachment():
    """别的工具返回里出现 attachmentId 可能只是引用了已有附件，
    重复挂到消息上会让同一个文件在历史里出现两次。"""
    from app.agents.conversation import _produced_attachment_id

    assert _produced_attachment_id("create_document", {"attachmentId": "a1"}) == "a1"
    assert _produced_attachment_id("site_resource_list", {"attachmentId": "a1"}) is None
    assert _produced_attachment_id("create_document", {}) is None
    assert _produced_attachment_id("create_document", "规格有问题：…") is None
    assert _produced_attachment_id("create_document", {"attachmentId": ""}) is None


async def test_generated_file_is_attached_to_the_assistant_message(
    authenticated_client,
    session_maker,
):
    """模型做完一份文档，历史里必须还能找到它——不能只活在工具返回值里。"""
    from types import SimpleNamespace

    from app.agents.conversation import AgentRuntime
    from app.api import get_agent_runtime
    from app.conversations import ConversationService

    class DocumentAgent:
        async def astream_events(self, *args, **kwargs):
            del args, kwargs
            yield {
                "event": "on_tool_start",
                "name": "create_document",
                "data": {"input": {"kind": "docx"}},
            }
            yield {
                "event": "on_tool_end",
                "name": "create_document",
                "data": {"output": '{"attachmentId":"att-123","filename":"周报.docx"}'},
            }
            yield {
                "event": "on_chat_model_stream",
                "data": {"chunk": SimpleNamespace(content="做好了")},
            }

    class FakeEmbedding:
        dimensions = 1024
        provider_name = "fake"
        model_name = "fake-1024"

        async def embed_query(self, text):
            del text
            return [0.1] * 1024

        async def embed_documents(self, texts):
            return [[0.1] * 1024 for _ in texts]

    runtime = AgentRuntime(DocumentAgent(), session_maker, FakeEmbedding())
    app = authenticated_client._transport.app
    app.dependency_overrides[get_agent_runtime] = lambda: runtime

    response = await authenticated_client.post(
        "/api/v1/chat/stream",
        json={"conversationId": None, "message": "帮我写个周报"},
    )
    # 流里立刻推一条，前端不用等整轮结束才看得见文件
    assert "event: attachment.ready" in response.text
    assert '"attachmentId":"att-123"' in response.text

    async with session_maker() as db:
        user_id = await db.scalar(select(User.id))
        conversation = (await ConversationService().list(db, user_id))[0]
        messages = await ConversationService().messages(db, user_id, conversation.id)
    assistant = [item for item in messages if item.role == "assistant"][-1]
    # 落库的那条也带着，刷新页面后还找得到
    assert assistant.metadata_["attachmentIds"] == ["att-123"]
    app.dependency_overrides.pop(get_agent_runtime, None)
