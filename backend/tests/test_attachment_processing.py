from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation

from app.attachment_processing import extract_text
from app.config import Settings
from app.document_services import render_preview_pdf, understand_document

LIMITS = {
    "max_chars": 20,
    "max_pdf_pages": 5,
    "max_office_uncompressed_bytes": 100,
    "max_workbook_sheets": 2,
    "max_workbook_rows": 10,
    "max_workbook_cells": 100,
}


def test_text_extraction_stops_at_character_limit() -> None:
    result = extract_text(
        "这是一个很长的附件内容，用来验证字符截断。".encode(),
        "text/plain",
        "note.txt",
        **LIMITS,
    )
    assert result is not None
    assert len(result) <= LIMITS["max_chars"]


def test_office_archive_rejects_excessive_expanded_size() -> None:
    archive_bytes = BytesIO()
    with ZipFile(archive_bytes, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "x" * 101)

    with pytest.raises(ValueError, match="解压后超过处理上限"):
        extract_text(
            archive_bytes.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "oversized.docx",
            **LIMITS,
        )


def test_pptx_fallback_reads_slides_and_notes() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "季度结论"
    slide.placeholders[1].text = "收入增长"
    slide.notes_slide.notes_text_frame.text = "提醒：不要逐字念"
    output = BytesIO()
    presentation.save(output)

    result = extract_text(
        output.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "report.pptx",
        **{**LIMITS, "max_chars": 200, "max_office_uncompressed_bytes": 1_000_000},
    )
    assert "季度结论" in result
    assert "收入增长" in result
    assert "不要逐字念" in result


async def test_builtin_document_understanding_is_explicitly_degraded() -> None:
    settings = Settings(session_secret="x" * 32, document_parser_url="")
    result = await understand_document(b"hello", "text/plain", "hello.txt", settings)
    assert result.text == "hello"
    assert result.parser == "builtin-degraded"
    assert result.processing_metadata["degraded"] is True


async def test_docling_adapter_preserves_text_and_structured_ir(monkeypatch) -> None:
    captured = {}

    class Response:
        def raise_for_status(self):
            return None

        def json(self):
            return {
                "status": "success",
                "document": {
                    "text_content": "结构化正文",
                    "md_content": "# 结构化正文",
                    "json_content": {"schema_name": "DoclingDocument", "pages": {"1": {}}},
                },
                "processing_time": 1.25,
                "timings": {"pipeline": 1.0},
                "errors": [],
            }

    class Client:
        async def __aenter__(self):
            return self

        async def __aexit__(self, *args):
            return None

        async def post(self, url, **kwargs):
            captured.update(url=url, **kwargs)
            return Response()

    monkeypatch.setattr("app.document_services.httpx.AsyncClient", lambda **kwargs: Client())
    settings = Settings(
        session_secret="x" * 32,
        document_parser_url="http://docling:5001",
        document_parser_api_key="internal-secret",
    )
    result = await understand_document(b"binary", "application/pdf", "sample.pdf", settings)
    assert captured["url"].endswith("/v1/convert/file")
    assert captured["headers"]["X-Api-Key"] == "internal-secret"
    assert captured["data"]["to_formats"] == ["md", "json", "text"]
    assert result.parser == "docling-v1"
    assert result.text == "结构化正文"
    assert result.document_ir["schema_name"] == "DoclingDocument"


async def test_gotenberg_adapter_returns_pdf_only_for_office(monkeypatch) -> None:
    class Response:
        content = b"%PDF-preview"
        headers = {"content-type": "application/pdf"}

        def raise_for_status(self):
            return None

    class Client:
        async def __aenter__(self):
            return self

        async def __aexit__(self, *args):
            return None

        async def post(self, url, **kwargs):
            assert url.endswith("/forms/libreoffice/convert")
            assert kwargs["data"]["exportNotes"] == "true"
            return Response()

    monkeypatch.setattr("app.document_services.httpx.AsyncClient", lambda **kwargs: Client())
    settings = Settings(session_secret="x" * 32, document_renderer_url="http://gotenberg:3000")
    preview = await render_preview_pdf(
        b"office",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "report.docx",
        settings,
    )
    assert preview == b"%PDF-preview"
