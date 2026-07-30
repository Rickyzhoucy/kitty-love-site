"""从结构化规格生成 docx / pptx / xlsx。

**为什么是工具而不是 Skill**：Skill 沙箱存在的理由是「执行模型或用户写的代码
不可信」。这里模型交出来的是一份 JSON 规格，代码是我们自己的——不需要沙箱，
也就不必付沙箱的代价（没有外网、只能靠 stdout 传结果、30 秒超时）。

反过来说，这也是**边界**：这个模块永远只按规格拼装文档，绝不 eval 规格里的
任何东西。想要「让模型写脚本处理数据」，那才是 Skill 该干的事。
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Any, Literal

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt

DocumentKind = Literal["docx", "pptx", "xlsx"]

CONTENT_TYPES: dict[str, str] = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

#: 规格的规模上限。模型偶尔会生成上千条内容，不设上限会把内存和生成时间打爆。
MAX_BLOCKS = 200
MAX_BULLETS_PER_BLOCK = 60
MAX_TEXT_CHARS = 5_000
MAX_TABLE_ROWS = 200
MAX_TABLE_COLS = 20


class DocumentSpecError(ValueError):
    """规格不合法。信息会原样回给模型，让它自己改。"""


@dataclass(frozen=True)
class BuiltDocument:
    filename: str
    content_type: str
    content: bytes


def _text(value: Any, limit: int = MAX_TEXT_CHARS) -> str:
    return str(value if value is not None else "").strip()[:limit]


def _rows(value: Any) -> list[list[str]]:
    if not isinstance(value, list):
        raise DocumentSpecError("表格的 rows 必须是二维数组")
    rows: list[list[str]] = []
    for row in value[:MAX_TABLE_ROWS]:
        if not isinstance(row, list):
            raise DocumentSpecError("表格的每一行必须是数组")
        rows.append([_text(cell, 500) for cell in row[:MAX_TABLE_COLS]])
    if not rows:
        raise DocumentSpecError("表格至少要有一行")
    return rows


def _bullets(value: Any) -> list[str]:
    if not isinstance(value, list):
        raise DocumentSpecError("bullets 必须是字符串数组")
    return [_text(item, 500) for item in value[:MAX_BULLETS_PER_BLOCK] if _text(item, 500)]


def _check_blocks(blocks: Any) -> list[dict[str, Any]]:
    if not isinstance(blocks, list) or not blocks:
        raise DocumentSpecError("blocks 必须是非空数组")
    if len(blocks) > MAX_BLOCKS:
        raise DocumentSpecError(f"内容块最多 {MAX_BLOCKS} 个")
    for block in blocks:
        if not isinstance(block, dict):
            raise DocumentSpecError("每个内容块必须是对象")
    return blocks


def build_docx(title: str, blocks: list[dict[str, Any]]) -> bytes:
    document = Document()
    if title:
        document.add_heading(title, level=0)
    for block in blocks:
        kind = str(block.get("type", "paragraph"))
        if kind == "heading":
            level = max(1, min(int(block.get("level", 1) or 1), 4))
            document.add_heading(_text(block.get("text"), 300), level=level)
        elif kind == "bullets":
            for item in _bullets(block.get("items")):
                document.add_paragraph(item, style="List Bullet")
        elif kind == "numbers":
            for item in _bullets(block.get("items")):
                document.add_paragraph(item, style="List Number")
        elif kind == "table":
            rows = _rows(block.get("rows"))
            table = document.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for row_index, row in enumerate(rows):
                for cell_index, cell in enumerate(row[: len(rows[0])]):
                    table.cell(row_index, cell_index).text = cell
        elif kind == "quote":
            paragraph = document.add_paragraph(_text(block.get("text")))
            paragraph.style = "Intense Quote"
        else:
            document.add_paragraph(_text(block.get("text")))
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pptx(title: str, blocks: list[dict[str, Any]]) -> bytes:
    presentation = Presentation()
    if title:
        cover = presentation.slides.add_slide(presentation.slide_layouts[0])
        cover.shapes.title.text = _text(title, 200)
        if len(cover.placeholders) > 1:
            cover.placeholders[1].text = _text(blocks[0].get("subtitle") or "", 200)

    for block in blocks:
        heading = _text(block.get("title") or block.get("text"), 200)
        bullets = _bullets(block.get("bullets") or block.get("items") or [])
        # 只带 subtitle 的块是封面副标题，上面已经用掉了。不跳过的话会多出
        # 一页彻底空白的幻灯片——模型很自然会把副标题写成第一个块。
        if not heading and not bullets:
            continue
        # 只有标题没有要点时用「节标题」版式，塞进正文版式会留一大块空白。
        layout = presentation.slide_layouts[1 if bullets else 5]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = heading
        if bullets:
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, item in enumerate(bullets):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = item
                paragraph.font.size = PptPt(18)
        notes = _text(block.get("notes"), 2_000)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    if not presentation.slides:
        raise DocumentSpecError("演示文稿至少要有一页")
    # 触发一次尺寸访问，确保模板正常加载（空演示文稿保存会失败得很隐晦）
    Inches(1)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_xlsx(title: str, blocks: list[dict[str, Any]]) -> bytes:
    workbook = Workbook()
    workbook.remove(workbook.active)
    for index, block in enumerate(blocks, start=1):
        name = _text(block.get("sheet") or block.get("title") or f"{title or 'Sheet'}{index}", 31)
        # Excel 的工作表名不能含这些字符，也不能为空
        for banned in "[]:*?/\\":
            name = name.replace(banned, "-")
        sheet = workbook.create_sheet(name or f"Sheet{index}")
        for row in _rows(block.get("rows")):
            sheet.append(row)
        if sheet.max_row >= 1:
            for cell in sheet[1]:
                cell.font = Font(bold=True)
    if not workbook.sheetnames:
        raise DocumentSpecError("工作簿至少要有一张表")
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _safe_filename(name: str, kind: str) -> str:
    stem = "".join(
        char for char in _text(name, 80) if char not in '\\/:*?"<>|\r\n\t'
    ).strip() or "未命名"
    return stem if stem.lower().endswith(f".{kind}") else f"{stem}.{kind}"


def build_document(
    kind: str,
    title: str,
    blocks: Any,
    filename: str | None = None,
) -> BuiltDocument:
    normalized = str(kind or "").lower().lstrip(".")
    if normalized not in CONTENT_TYPES:
        raise DocumentSpecError(
            f"不支持的文档类型 {kind}，可选：{', '.join(sorted(CONTENT_TYPES))}"
        )
    checked = _check_blocks(blocks)
    clean_title = _text(title, 200)
    builders = {"docx": build_docx, "pptx": build_pptx, "xlsx": build_xlsx}
    content = builders[normalized](clean_title, checked)
    return BuiltDocument(
        filename=_safe_filename(filename or clean_title, normalized),
        content_type=CONTENT_TYPES[normalized],
        content=content,
    )
