from io import BytesIO
from zipfile import BadZipFile, ZipFile

from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps
from pptx import Presentation
from pypdf import PdfReader


def _bounded_append(parts: list[str], value: str, current: int, limit: int) -> int:
    remaining = limit - current
    if remaining <= 0:
        return current
    piece = value[:remaining]
    parts.append(piece)
    return current + len(piece)


def _validate_office_archive(content: bytes, expanded_limit: int) -> None:
    try:
        with ZipFile(BytesIO(content)) as archive:
            expanded = sum(entry.file_size for entry in archive.infolist())
    except BadZipFile as error:
        raise ValueError("Office 文件不是有效的 ZIP 容器") from error
    if expanded > expanded_limit:
        raise ValueError("Office 文件解压后超过处理上限")


def extract_text(
    content: bytes,
    content_type: str,
    filename: str,
    *,
    max_chars: int,
    max_pdf_pages: int,
    max_office_uncompressed_bytes: int,
    max_workbook_sheets: int,
    max_workbook_rows: int,
    max_workbook_cells: int,
) -> str | None:
    suffix = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if content_type.startswith("text/") or suffix in {"txt", "md", "csv", "json"}:
        return content[: max_chars * 4].decode("utf-8", errors="replace")[:max_chars]
    if content_type == "application/pdf" or suffix == "pdf":
        reader = PdfReader(BytesIO(content))
        if len(reader.pages) > max_pdf_pages:
            raise ValueError("PDF 页数超过处理上限")
        parts: list[str] = []
        length = 0
        for page in reader.pages:
            length = _bounded_append(parts, page.extract_text() or "", length, max_chars)
            if length >= max_chars:
                break
            length = _bounded_append(parts, "\n\n", length, max_chars)
        return "".join(parts)
    if suffix == "docx":
        _validate_office_archive(content, max_office_uncompressed_bytes)
        document = Document(BytesIO(content))
        parts = []
        length = 0
        for paragraph in document.paragraphs:
            length = _bounded_append(parts, paragraph.text, length, max_chars)
            if length >= max_chars:
                break
            length = _bounded_append(parts, "\n", length, max_chars)
        for table_index, table in enumerate(document.tables, start=1):
            length = _bounded_append(parts, f"\n## 表格 {table_index}\n", length, max_chars)
            for row in table.rows:
                text = "\t".join(cell.text for cell in row.cells)
                length = _bounded_append(parts, f"{text}\n", length, max_chars)
                if length >= max_chars:
                    break
            if length >= max_chars:
                break
        return "".join(parts)
    if suffix == "xlsx":
        _validate_office_archive(content, max_office_uncompressed_bytes)
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        parts = []
        length = 0
        row_count = 0
        cell_count = 0
        if len(workbook.worksheets) > max_workbook_sheets:
            workbook.close()
            raise ValueError("工作表数量超过处理上限")
        for sheet in workbook.worksheets[:max_workbook_sheets]:
            length = _bounded_append(parts, f"# {sheet.title}\n", length, max_chars)
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                cell_count += len(row)
                if (
                    row_count > max_workbook_rows
                    or cell_count > max_workbook_cells
                ):
                    workbook.close()
                    raise ValueError("工作簿行数或单元格数量超过处理上限")
                text = "\t".join("" if value is None else str(value) for value in row)
                length = _bounded_append(parts, f"{text}\n", length, max_chars)
                if length >= max_chars:
                    break
            if length >= max_chars:
                break
        workbook.close()
        return "".join(parts)
    if suffix == "pptx":
        _validate_office_archive(content, max_office_uncompressed_bytes)
        presentation = Presentation(BytesIO(content))
        parts = []
        length = 0
        for index, slide in enumerate(presentation.slides, start=1):
            length = _bounded_append(parts, f"# 幻灯片 {index}\n", length, max_chars)
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    length = _bounded_append(parts, f"{text}\n", length, max_chars)
            notes = slide.notes_slide
            for shape in notes.shapes:
                if getattr(shape, "is_placeholder", False):
                    text = getattr(shape, "text", "")
                    if text and text.strip() not in {str(index), ""}:
                        length = _bounded_append(parts, f"备注：{text}\n", length, max_chars)
            if length >= max_chars:
                break
        return "".join(parts)
    return None


def thumbnail_webp(content: bytes, content_type: str) -> bytes | None:
    if not content_type.startswith("image/"):
        return None
    with Image.open(BytesIO(content)) as image:
        # 手机照片常把方向只写进 EXIF；不先转正，缩略图会横躺或倒置，而原图在
        # 浏览器里又是正的，点开前后方向就跳了。
        image = ImageOps.exif_transpose(image)
        image.thumbnail((512, 512), Image.Resampling.LANCZOS)
        output = BytesIO()
        image.convert("RGB").save(output, "WEBP", quality=82, method=6)
        return output.getvalue()
